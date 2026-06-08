"""Tests for the PPTX renderer. Round-tripping through python-pptx is also our
proxy for 'opens cleanly': a file python-pptx re-parses is valid OOXML."""

from pptx import Presentation

from polymath.models.schemas import Slide, SlideDeck
from polymath.outputs.slides import render_deck


def _deck() -> SlideDeck:
    return SlideDeck(
        title="Solid-State Batteries",
        slides=[
            Slide(title="Safety", bullets=["non-flammable", "no thermal runaway"], source_url="https://example.com/a"),
            Slide(title="Timeline", bullets=["2027 pilot", "2030 mass production"], source_url="https://example.com/b"),
        ],
    )


def test_render_creates_openable_pptx(tmp_path):
    out = render_deck(_deck(), tmp_path / "deck.pptx")
    assert out.exists()
    prs = Presentation(str(out))  # re-parse == valid OOXML
    assert len(prs.slides) == 3  # 1 title + 2 content


def test_title_and_finding_text_present(tmp_path):
    out = render_deck(_deck(), tmp_path / "deck.pptx")
    prs = Presentation(str(out))
    all_text = "\n".join(
        shape.text for slide in prs.slides for shape in slide.shapes if shape.has_text_frame
    )
    assert "Solid-State Batteries" in all_text
    assert "Safety" in all_text
    assert "non-flammable" in all_text
    assert "https://example.com/a" in all_text


def test_bullets_capped_at_three(tmp_path):
    deck = SlideDeck(
        title="T",
        slides=[Slide(title="Many", bullets=["b1", "b2", "b3", "b4", "b5"], source_url="https://x")],
    )
    out = render_deck(deck, tmp_path / "deck.pptx")
    prs = Presentation(str(out))
    body_text = "\n".join(
        shape.text for slide in prs.slides for shape in slide.shapes if shape.has_text_frame
    )
    assert "b3" in body_text
    assert "b4" not in body_text  # capped to 3 bullets
